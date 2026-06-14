"""Build a professional PowerPoint deck for the weekly research update."""

from __future__ import annotations

import json
import os
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# Paths
RESULTS_FILE = Path("results/unified_evaluation.json")
OUTPUT_DIR = Path("results/slides")
OUTPUT_DIR.mkdir(exist_ok=True)

# Design constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)

# Color palette (professional, muted)
DARK_BLUE = RGBColor(31, 78, 121)
ACCENT = RGBColor(68, 114, 148)
LIGHT_BLUE = RGBColor(189, 215, 238)
ORANGE = RGBColor(237, 125, 49)
GRAY = RGBColor(127, 127, 127)
DARK_GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def load_results():
    with open(RESULTS_FILE) as f:
        return json.load(f)


def make_accuracy_chart(data, path):
    methods = list(data.keys())
    accs = [data[m]["accuracy"] * 100 for m in methods]

    colors = [ORANGE if m == "ACC" else DARK_BLUE for m in methods]

    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars = ax.bar(methods, accs, color=[(c[0]/255, c[1]/255, c[2]/255) for c in colors], edgecolor="white", linewidth=1.2)

    # Add value labels on bars
    for bar, val in zip(bars, accs):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.5,
                f"{val:.0f}%", ha="center", va="bottom", fontsize=14, fontweight="bold", color=(64/255, 64/255, 64/255))

    ax.axhline(y=50, color=(153/255, 153/255, 153/255), linestyle="--", linewidth=1.5, label="Random chance")
    ax.set_ylim(0, 100)
    ax.set_ylabel("Accuracy (%)", fontsize=13, color=(64/255, 64/255, 64/255))
    ax.set_title("Generation Accuracy by Method", fontsize=16, fontweight="bold", color=(31/255, 78/255, 121/255), pad=15)
    ax.tick_params(axis="x", labelsize=12, colors="#404040")
    ax.tick_params(axis="y", labelsize=11, colors="#404040")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(fontsize=11, frameon=False)

    plt.tight_layout()
    plt.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()


def make_detection_chart(data, path):
    methods = [m for m in data.keys() if m != "Baseline"]
    f1s = [data[m]["f1"] for m in methods]
    flag_rates = [data[m]["flag_rate"] * 100 for m in methods]

    x = np.arange(len(methods))
    width = 0.35

    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars1 = ax.bar(x - width / 2, f1s, width, label="Detection F1", color=(31/255, 78/255, 121/255), edgecolor="white")
    bars2 = ax.bar(x + width / 2, flag_rates, width, label="Flag Rate (%)", color=(127/255, 127/255, 127/255), edgecolor="white")

    # Labels
    for bar, val in zip(bars1, f1s):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.02,
                f"{val:.2f}", ha="center", va="bottom", fontsize=11, fontweight="bold", color=(64/255, 64/255, 64/255))
    for bar, val in zip(bars2, flag_rates):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                f"{val:.0f}%", ha="center", va="bottom", fontsize=11, fontweight="bold", color=(64/255, 64/255, 64/255))

    ax.set_ylim(0, 1.2)
    ax.set_ylabel("Score", fontsize=13, color=(64/255, 64/255, 64/255))
    ax.set_title("Detection Quality vs. Selectivity", fontsize=16, fontweight="bold", color=(31/255, 78/255, 121/255), pad=15)
    ax.set_xticks(x)
    ax.set_xticklabels(methods, fontsize=12)
    ax.tick_params(axis="y", labelsize=11, colors="#404040")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(fontsize=11, frameon=False, loc="upper right")

    plt.tight_layout()
    plt.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()


def make_per_type_chart(path):
    # Per-type data from the weekly update
    categories = ["Factual", "Hallucination", "Uncertain"]
    baseline = [50, 50, 50]
    acc = [75, 50, 50]

    x = np.arange(len(categories))
    width = 0.35

    fig, ax = plt.subplots(figsize=(7, 4.2))
    bars1 = ax.bar(x - width / 2, baseline, width, label="Baseline", color=(127/255, 127/255, 127/255), edgecolor="white")
    bars2 = ax.bar(x + width / 2, acc, width, label="ACC", color=(237/255, 125/255, 49/255), edgecolor="white")

    for bar, val in zip(bars1, baseline):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.5,
                f"{val}%", ha="center", va="bottom", fontsize=11, fontweight="bold", color=(64/255, 64/255, 64/255))
    for bar, val in zip(bars2, acc):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.5,
                f"{val}%", ha="center", va="bottom", fontsize=11, fontweight="bold", color=(64/255, 64/255, 64/255))

    ax.set_ylim(0, 100)
    ax.set_ylabel("Accuracy (%)", fontsize=13, color=(64/255, 64/255, 64/255))
    ax.set_title("ACC Improvement is Driven by Factual Questions", fontsize=15, fontweight="bold", color=(31/255, 78/255, 121/255), pad=15)
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=12)
    ax.tick_params(axis="y", labelsize=11, colors="#404040")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(fontsize=11, frameon=False)

    plt.tight_layout()
    plt.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close()


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def add_bullet_box(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
    return box


def add_title_bar(slide, title, subtitle=None):
    # Top color bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    add_text_box(slide, MARGIN, Inches(0.22), Inches(12), Inches(0.7),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, MARGIN, Inches(0.82), Inches(12), Inches(0.4),
                     subtitle, font_size=14, color=LIGHT_BLUE)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    data = load_results()

    # Generate charts
    chart1 = OUTPUT_DIR / "accuracy_chart.png"
    chart2 = OUTPUT_DIR / "detection_chart.png"
    chart3 = OUTPUT_DIR / "per_type_chart.png"
    make_accuracy_chart(data, chart1)
    make_detection_chart(data, chart2)
    make_per_type_chart(chart3)

    # Ablation chart (re-use existing if available)
    ablation_chart = OUTPUT_DIR / "ablation_chart.png"
    if not ablation_chart.exists() and Path("results/figures/ablation_accuracy.png").exists():
        from PIL import Image
        img = Image.open("results/figures/ablation_accuracy.png")
        img.save(ablation_chart)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_text_box(slide, MARGIN, Inches(2.2), Inches(12), Inches(1.5),
                 "ACC: Neuroscience-Inspired Hallucination Detection",
                 font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide, MARGIN, Inches(3.6), Inches(12), Inches(0.8),
                 "Weekly Research Update — June 28, 2026",
                 font_size=22, color=DARK_GRAY)
    add_text_box(slide, MARGIN, Inches(4.5), Inches(12), Inches(0.6),
                 "Hafez Al-Khatib",
                 font_size=18, color=GRAY)
    # Bottom bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Problem: LLMs Hallucinate Without Self-Correction")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Large language models generate fluent but ungrounded statements.",
                       "Existing methods detect hallucination after output is complete (post-hoc).",
                       "Our goal: intervene during generation, inspired by conflict monitoring in the anterior cingulate cortex (ACC).",
                       "Core hypothesis: prediction errors between hierarchical layers signal when the model is departing from reliable patterns.",
                   ], font_size=18)

    # Slide 3: Approach
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Approach: Predictive Coding + Conflict Monitoring")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(6), Inches(5),
                   [
                       "Predictive coding: each layer predicts the next layer's representation.",
                       "Large prediction error = potential conflict.",
                       "Leaky integrator accumulates evidence over tokens.",
                       "Decision engine maps conflict + entropy to actions: flag, regenerate, warn.",
                   ], font_size=17)
    # Simple architecture diagram via text
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(4.5),
                 "Generation → Hidden States\n↓\nPrediction-Error Features\n↓\nMLP Detector\n↓\nConflict Score + Entropy\n↓\nDecision Engine\n↓\nIntervention",
                 font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Slide 4: Pipeline built this week
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What We Built This Week")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Baseline detectors: DoLa, SAPLMA, Entropy.",
                       "ACC intervention engine: draft → detect → regenerate with uncertainty priming.",
                       "Model-specific detector training on Qwen2.5-1.5B (600 token-level examples).",
                       "Unified evaluation script with fair comparison controls.",
                       "Colab GPU notebook for larger-scale experiments.",
                   ], font_size=18)

    # Slide 5: Evaluation methodology
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Evaluation: Fair Comparison Controls")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Same 10 prompts across all methods.",
                       "Same random seeds for reproducibility.",
                       "Same softmax sampling strategy (no unfair top-p differences).",
                       "Same judge function: substring-based correctness + uncertainty markers.",
                       "SAPLMA trained on held-out prompts to avoid data leakage.",
                   ], font_size=18)

    # Slide 6: Accuracy results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Results: Generation Accuracy")
    slide.shapes.add_picture(str(chart1), MARGIN, Inches(1.4), width=Inches(12))
    add_text_box(slide, MARGIN, Inches(6.5), Inches(12), Inches(0.5),
                 "ACC improves over baseline, but absolute numbers remain modest.",
                 font_size=14, color=GRAY)

    # Slide 7: Detection results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Results: Detection Quality vs. Selectivity")
    slide.shapes.add_picture(str(chart2), MARGIN, Inches(1.4), width=Inches(12))
    add_text_box(slide, MARGIN, Inches(6.5), Inches(12), Inches(0.5),
                 "DoLa has highest F1 but flags 100% of tokens. ACC balances F1 and selectivity.",
                 font_size=14, color=GRAY)

    # Slide 8: Per-type
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Results: ACC Improvement by Question Type")
    slide.shapes.add_picture(str(chart3), Inches(3.5), Inches(1.5), width=Inches(6.5))
    add_text_box(slide, MARGIN, Inches(6.2), Inches(12), Inches(0.7),
                 "ACC gains come from factual questions. Hallucination and uncertainty detection remain at chance on this model.",
                 font_size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Slide 9: Logit-shift intervention
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Stronger Intervention: Logit-Shift")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Phrase-based intervention only changes the prompt prefix.",
                       "Logit-shift directly biases the output distribution toward uncertainty tokens.",
                       "Implemented in src/acc_intervention.py: generate_with_logit_shift().",
                       "Tokens like 'I am not sure', 'actually', 'wait' receive a positive bias when conflict is high.",
                   ], font_size=18)

    # Slide 10: Ablation results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Ablation Study: Component Contributions")
    if ablation_chart.exists():
        slide.shapes.add_picture(str(ablation_chart), MARGIN, Inches(1.4), width=Inches(12))
    else:
        add_text_box(slide, MARGIN, Inches(2.5), Inches(12), Inches(2),
                     "Run scripts/run_ablation.py to generate this chart.",
                     font_size=18, color=GRAY)
    add_text_box(slide, MARGIN, Inches(6.5), Inches(12), Inches(0.5),
                 "Logit-shift (absolute threshold) reaches 80% vs. 50% baseline on 1.5B.",
                 font_size=14, color=GRAY)

    # Slide 11: Key insight
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Insight: Scale Is the Bottleneck")
    add_text_box(slide, MARGIN, Inches(2.0), Inches(12), Inches(2.0),
                 "The detector works (83% validation accuracy), and the intervention pipeline works, but Qwen2.5-1.5B lacks the capacity to consistently act on the signal.",
                 font_size=22, color=DARK_BLUE, align=PP_ALIGN.LEFT)
    add_bullet_box(slide, MARGIN, Inches(4.1), Inches(12), Inches(2.5),
                   [
                       "Factual questions improve because the model knows the answer.",
                       "Hallucination prompts produce inconsistent behavior regardless of detector.",
                       "Next step: evaluate on Qwen2.5-7B via Colab GPU.",
                   ], font_size=17)

    # Slide 12: Bugs fixed
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Critical Bugs Found & Fixed")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Inverted synthetic logits: detector always predicted 'supported'. Fixed mapping.",
                       "Post-hoc mismatch: detection re-ran generation with different random sample. Now caches hidden states during main pass.",
                       "Gradient leak: explainability backward() accumulated gradients on model weights. Now uses torch.autograd.grad.",
                       "Plus: unfair sampling, SAPLMA data leakage, missing hook cleanup.",
                   ], font_size=17)

    # Slide 13: Audit
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Code Audit: Quality Now Validated")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "Audited 7 core files across architecture, integration, baselines, evaluation.",
                       "Identified 20 ranked issues in results/code_audit_report.md.",
                       "All critical/high issues resolved or in progress.",
                       "Logit-shift intervention and ablation study completed this week.",
                   ], font_size=18)

    # Slide 14: Next step
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Next Step: Large-Scale Benchmark Evaluation")
    add_bullet_box(slide, MARGIN, Inches(1.5), Inches(12), Inches(5),
                   [
                       "scripts/run_benchmark_eval.py: HaluEval + TruthfulQA + PubMedQA.",
                       "LLM-as-judge, SelfCheckGPT baseline, bootstrap CIs, paired t-tests.",
                       "Run on Qwen2.5-7B via 4090 or Colab T4.",
                       "Target: 200+ samples per benchmark for publication credibility.",
                   ], font_size=18)

    # Slide 15: Plan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Next Week's Plan")
    table = slide.shapes.add_table(5, 3, MARGIN, Inches(1.5), Inches(12), Inches(5)).table
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(6.5)
    table.columns[2].width = Inches(4.7)

    rows = [
        ("1", "Run benchmark evaluation on Qwen2.5-7B (200+ samples)", "Publication-scale metrics"),
        ("2", "Run full ablation on 7B (detector, threshold, intervention)", "Component importance"),
        ("3", "Add more baselines (SelfCheckGPT, LMvsLM)", "Stronger comparison"),
        ("4", "Draft methods section with neuroscience framework", "Paper structure"),
        ("5", "Generate final figures for AAAI submission", "Camera-ready visuals"),
    ]
    for i, (p, task, out) in enumerate(rows):
        table.cell(i, 0).text = p
        table.cell(i, 1).text = task
        table.cell(i, 2).text = out
        for c in range(3):
            cell = table.cell(i, c)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else WHITE

    # Slide 16: Discussion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Discussion Questions")
    add_bullet_box(slide, MARGIN, Inches(1.8), Inches(12), Inches(4.5),
                   [
                       "Can we secure 4090 access this week for 7B benchmark runs?",
                       "Should we invest in paid GPU hours if local/Colab access is unstable?",
                       "Target venue: AAAI workshop first, then top-tier after scale + baselines?",
                   ], font_size=22)

    # Save
    output_path = OUTPUT_DIR / "ACC_Weekly_Update_2026-06-28.pptx"
    prs.save(output_path)
    print(f"Saved slides to: {output_path}")
    print(f"Charts saved to: {OUTPUT_DIR}")


if __name__ == "__main__":
    build_deck()
